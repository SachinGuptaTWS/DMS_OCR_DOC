import os
import hashlib
import concurrent.futures
import time
import json
import warnings
import io

# --- 1. Library Imports ---
# Ensure these are installed: pip install python-pptx easyocr numpy opencv-python-headless pillow gliner google-generativeai
try:
    import pptx
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import easyocr
    import numpy as np
    import cv2
    from PIL import Image
    from gliner import GLiNER # For "Powerful" Entity Extraction
    import google.generativeai as genai # For Vision AI
except ImportError as e:
    print(f"CRITICAL: Missing Dependency in _ppt_ocr.py: {e}")

# Suppress warnings
warnings.filterwarnings("ignore")

# API Configuration
GEMINI_API_KEY = "AIzaSyCI3S7weovcMknn3sSpsTstB9N4b57VvqY"

# --- 2. Gemini Vision Engine (NEW) ---
class GeminiVisionEngine:
    def __init__(self, api_key):
        self.available = False
        try:
            if not api_key:
                print("WARNING: No Gemini API Key provided.")
                return
            genai.configure(api_key=api_key)
            self.model = genai.GenerativeModel('gemini-1.5-flash')
            self.available = True
            print("INFO: Gemini Vision AI Initialized.")
        except Exception as e:
            print(f"WARNING: Gemini Init Failed: {e}")

    def analyze_image(self, blob):
        """
        Sends image blob to Gemini 1.5 Flash for advanced transcription.
        """
        if not self.available: return None
        try:
            # Convert bytes to PIL Image
            image = Image.open(io.BytesIO(blob))
            
            # Prompt optimized for extraction
            prompt = (
                "Analyze this image. "
                "1. If it contains text, transcribe it exactly. "
                "2. If it contains a table, format it as a Markdown table. "
                "3. If it contains a chart or diagram, provide a concise summary of the data/insight. "
                "Return only the content, no conversational filler."
            )
            
            response = self.model.generate_content([prompt, image])
            return response.text.strip()
        except Exception as e:
            print(f"Gemini Inference Error: {e}")
            return None

# --- 3. Image Enhancer (Pre-processing for OCR) ---
class ImageEnhancer:
    @staticmethod
    def process_image_variants(blob):
        variants = []
        try:
            nparr = np.frombuffer(blob, np.uint8)
            original = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
            if original is None: return []

            # 1. Grayscale + Denoise (General purpose)
            gray = cv2.cvtColor(original, cv2.COLOR_BGR2GRAY)
            denoised = cv2.fastNlMeansDenoising(gray, h=10, templateWindowSize=7, searchWindowSize=21)
            variants.append(denoised)

            # 2. CLAHE (Contrast Limited Adaptive Histogram Equalization) - Good for dark/faded text
            clahe = cv2.createCLAHE(clipLimit=3.0, tileGridSize=(8,8))
            contrast = clahe.apply(gray)
            variants.append(contrast)

            # 3. Binarization (Otsu) - Good for high contrast documents
            # Upscale small images to help OCR
            h, w = gray.shape
            scale = 2 if w < 1000 else 1
            if scale > 1:
                resized = cv2.resize(gray, (w*scale, h*scale), interpolation=cv2.INTER_CUBIC)
            else:
                resized = gray
            _, binary = cv2.threshold(resized, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            variants.append(binary)

        except Exception:
            pass 
        return variants

# --- 4. OCR Engine (EasyOCR + Gemini Hybrid) ---
class PowerOCREngine:
    def __init__(self, use_gpu=True, gemini_engine=None):
        print("INFO: Initializing EasyOCR...")
        try:
            self.reader = easyocr.Reader(['en'], gpu=use_gpu)
        except Exception as e:
            print(f"WARNING: GPU init failed ({e}). Fallback to CPU.")
            self.reader = easyocr.Reader(['en'], gpu=False)
            
        self.gemini = gemini_engine
        self.blob_cache = {} 
        self.queue = [] 

    def queue_image(self, blob, slide_idx):
        if not blob: return
        # Hash image to avoid re-processing duplicates (common in logos/footers)
        img_hash = hashlib.md5(blob).hexdigest()
        if img_hash not in self.blob_cache:
            self.queue.append({
                "hash": img_hash,
                "blob": blob,
                "slide_idx": slide_idx
            })

    def run_batch_process(self):
        unique_tasks = {item["hash"]: item["blob"] for item in self.queue}
        if not unique_tasks: return

        print(f"INFO: Processing {len(unique_tasks)} unique images...")

        # --- A. Local EasyOCR (Baseline) ---
        # 1. Pre-process
        image_variants_map = {}
        with concurrent.futures.ThreadPoolExecutor() as executor:
            future_to_hash = {executor.submit(ImageEnhancer.process_image_variants, blob): h 
                              for h, blob in unique_tasks.items()}
            for future in concurrent.futures.as_completed(future_to_hash):
                try:
                    variants = future.result()
                    if variants: image_variants_map[future_to_hash[future]] = variants
                except Exception: pass

        # 2. Batch Inference
        batch_input = []
        batch_mapping = []
        for h, variants in image_variants_map.items():
            for v in variants:
                batch_input.append(v)
                batch_mapping.append(h)

        if batch_input:
            results = self.reader.readtext_batched(batch_input, detail=1, batch_size=4, paragraph=False)
            
            # 3. Aggregation
            temp_results = {}
            for i, result in enumerate(results):
                text_items = [item[1] for item in result]
                conf_items = [item[2] for item in result]
                
                if text_items:
                    text_found = " ".join(text_items)
                    avg_conf = sum(conf_items) / len(conf_items) if conf_items else 0.0
                    
                    if len(text_found.strip()) > 2:
                        h = batch_mapping[i]
                        if h not in temp_results: temp_results[h] = []
                        temp_results[h].append((text_found, avg_conf))

            for h, candidates in temp_results.items():
                best_candidate = max(candidates, key=lambda x: len(x[0]))
                self.blob_cache[h] = {"text": best_candidate[0], "conf": best_candidate[1]}

        # --- B. Cloud Gemini Vision (Enhancement) ---
        if self.gemini and self.gemini.available:
            print("INFO: Enhancing images with Gemini Vision AI...")
            count = 0
            # Process sequentially to avoid aggressive rate limits on free tier
            for h, blob in unique_tasks.items():
                # Only use Gemini if EasyOCR failed or gave low confidence? 
                # For "Powerful" mode, we try Gemini on everything.
                gemini_text = self.gemini.analyze_image(blob)
                
                if gemini_text:
                    count += 1
                    # Overwrite local OCR result with Gemini's superior understanding
                    self.blob_cache[h] = {
                        "text": gemini_text, 
                        "conf": 0.99, # Artificial high confidence for AI
                        "source": "Gemini Vision"
                    }
                    time.sleep(1) # Gentle rate limiting
            print(f"INFO: Gemini processed {count} images successfully.")

# --- 5. NER Engine (GLiNER - EXPANDED ONTOLOGY) ---
class NEREngine:
    def __init__(self):
        print("INFO: Initializing GLiNER...")
        try:
            self.model = GLiNER.from_pretrained("urchade/gliner_small-v2.1")
            self.available = True
        except Exception as e:
            print(f"WARNING: GLiNER load failed: {e}")
            self.available = False

    def extract_entities(self, text):
        if not self.available or not text: return []
        
        # --- POWERFUL: Extensive List of Entity Types ---
        # GLiNER is a zero-shot model, so we can ask for anything we want.
        labels = [
            # People & Orgs
            "person", "organization", "company", "university", "department", "job title",
            
            # Geography
            "location", "city", "country", "address",
            
            # Technology & Product
            "technology", "software", "hardware", "product", "service", "programming language", 
            "database", "cloud service", "algorithm",
            
            # Business & Finance
            "project", "initiative", "strategy", "brand", "competitor",
            "money", "currency", "revenue", "cost", "budget", "profit", "stock ticker",
            
            # Metrics & Data
            "metric", "percentage", "quantity", "KPI", "date", "time", "duration",
            
            # Digital & Legal
            "email", "url", "ip address", "law", "regulation", "contract", "clause"
        ]
        
        try:
            # Predict with a slightly lower threshold to catch more variety
            entities = self.model.predict_entities(text, labels, threshold=0.25)
            
            # Return cleaned list
            return [{"text": e["text"], "label": e["label"], "score": round(e["score"], 2)} for e in entities]
        except Exception:
            return []

# --- 6. Main Parser (PPTX Structure) ---
class PresentationParser:
    def __init__(self, file_path, ocr_engine, ner_engine):
        self.file_path = file_path
        self.ocr = ocr_engine
        self.ner = ner_engine
        self.slides_data = []

    def _extract_table(self, shape):
        if not shape.has_table: return None
        rows = []
        for row in shape.table.rows:
            cell_texts = []
            for cell in row.cells:
                if cell.text_frame:
                    cell_text = cell.text_frame.text.strip().replace("\n", " ")
                    cell_texts.append(cell_text)
                else:
                    cell_texts.append("")
            rows.append(cell_texts)
        if not rows: return None
        
        md_lines = []
        if len(rows) > 0:
            md_lines.append("| " + " | ".join(rows[0]) + " |")
            md_lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
            for r in rows[1:]:
                md_lines.append("| " + " | ".join(r) + " |")
        return "\n".join(md_lines)

    def _extract_text_frame(self, shape):
        items = []
        if not shape.has_text_frame: return items
        text_lines = []
        for p in shape.text_frame.paragraphs:
            raw_text = p.text.strip()
            if raw_text:
                prefix = "- " if p.level == 0 else "  " * p.level + "- "
                text_lines.append(prefix + raw_text)
        full_text = "\n".join(text_lines)
        if full_text:
            items.append({"text": full_text, "conf": 1.0, "type": "text_block"})
        return items

    def _process_shape(self, shape, slide_idx):
        items = []
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                items.extend(self._process_shape(sub_shape, slide_idx))
            return items

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and getattr(shape, 'placeholder_format', None) and shape.placeholder_format.type == 18):
            if hasattr(shape, "image"):
                self.ocr.queue_image(shape.image.blob, slide_idx)

        if shape.has_table:
            table_md = self._extract_table(shape)
            if table_md:
                items.append({"text": table_md, "conf": 1.0, "type": "table_structure"})

        if shape.has_text_frame:
            items.extend(self._extract_text_frame(shape))

        return items

    def parse(self):
        if not os.path.exists(self.file_path): return [], []
        prs = pptx.Presentation(self.file_path)
        
        for i, slide in enumerate(prs.slides):
            slide_idx = i
            slide_items = []
            for shape in slide.shapes:
                slide_items.extend(self._process_shape(shape, slide_idx))
            self.slides_data.append({"id": slide_idx + 1, "items": slide_items})

        self.ocr.run_batch_process()

        final_slides_output = []
        all_text_accumulated = ""

        for slide in self.slides_data:
            slide_idx = slide["id"] - 1
            relevant_tasks = [t for t in self.ocr.queue if t["slide_idx"] == slide_idx]
            for task in relevant_tasks:
                h = task["hash"]
                if h in self.ocr.blob_cache:
                    data = self.ocr.blob_cache[h]
                    slide["items"].append({
                        "text": f"[IMAGE TEXT]: {data['text']}", 
                        "conf": data["conf"], 
                        "type": "ocr_image"
                    })

            full_text = "\n\n".join([item["text"] for item in slide["items"]])
            all_text_accumulated += f"\n--- Slide {slide['id']} ---\n{full_text}"
            
            if slide["items"]:
                avg_conf = sum(item["conf"] for item in slide["items"]) / len(slide["items"])
            else:
                avg_conf = 1.0

            final_slides_output.append({
                "slide_number": slide["id"],
                "text_content": full_text,
                "accuracy_score": round(avg_conf * 100, 2),
                "item_count": len(slide["items"])
            })

        print("INFO: Extracting Entities...")
        entities = self.ner.extract_entities(all_text_accumulated)
        return final_slides_output, entities

# --- 7. Main Entry Point ---
def execute_ppt_ocr(filepath):
    start_time = time.time()
    tool_name = "EasyOCR + GLiNER + Gemini Vision"
    
    try:
        import pptx, easyocr, gliner, google.generativeai
    except ImportError as e:
        return 0, {}, tool_name, f"Error: Missing libraries: {e}"

    try:
        gemini_engine = GeminiVisionEngine(GEMINI_API_KEY)
        ocr_engine = PowerOCREngine(use_gpu=True, gemini_engine=gemini_engine)
        ner_engine = NEREngine()
        
        parser = PresentationParser(filepath, ocr_engine, ner_engine)
        slides, entities = parser.parse()
        
        full_text_all = "\n\n".join([s["text_content"] for s in slides])
        total_acc = sum(s["accuracy_score"] for s in slides) / len(slides) if slides else 0.0

        detailed_result = {
            "total_accuracy": round(total_acc, 2),
            "total_slides": len(slides),
            "extracted_text": full_text_all,
            "per_slide_metrics": slides,
            "ner_entities": entities
        }
        
        duration = time.time() - start_time
        return duration, detailed_result, tool_name, "Success"
        
    except Exception as e:
        print(f"PPT Processing Error: {e}")
        return time.time() - start_time, {"error": str(e)}, tool_name, f"Error: {str(e)}"